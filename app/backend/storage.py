from pathlib import Path
from uuid import uuid4

from pptx import Presentation
from pypdf import PdfReader


PROJECT_ROOT = Path(__file__).resolve().parents[2]
DATA_DIR = PROJECT_ROOT / "data"
UPLOAD_DIR = DATA_DIR / "uploads"
CHROMA_DIR = DATA_DIR / "chroma"


def ensure_data_dirs() -> None:
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    CHROMA_DIR.mkdir(parents=True, exist_ok=True)


def save_upload(filename: str, content: bytes) -> Path:
    ensure_data_dirs()
    safe_name = Path(filename).name.replace(" ", "_")
    path = UPLOAD_DIR / f"{uuid4().hex}_{safe_name}"
    path.write_bytes(content)
    return path


def load_text(path: Path) -> list[dict]:
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        reader = PdfReader(str(path))
        pages = []
        for index, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append({"text": text, "page": index})
        return pages

    if suffix in {".txt", ".md"}:
        return [{"text": path.read_text(encoding="utf-8", errors="ignore"), "page": None}]

    if suffix == ".pptx":
        presentation = Presentation(str(path))
        slides = []
        for index, slide in enumerate(presentation.slides, start=1):
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
            slide_text = "\n".join(text_parts)
            if slide_text.strip():
                slides.append({"text": slide_text, "page": index})
        return slides

    raise ValueError("Only PDF, PPTX, TXT, and Markdown files are supported.")
